"""Tests for the sandbox report toolkit and its workspace seeding."""

import importlib.util
import sys
from pathlib import Path

import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
import pandas as pd
import pytest

REPO_ROOT = Path(__file__).resolve().parents[2]
KIT_PATH = REPO_ROOT / "sandbox" / "sandboxd" / "reportkit.py"
SERVER_PATH = REPO_ROOT / "sandbox" / "sandboxd" / "server.py"


def _load(name: str, path: Path):
    spec = importlib.util.spec_from_file_location(name, path)
    if spec is None or spec.loader is None:
        raise ImportError(f"cannot load {path}")
    module = importlib.util.module_from_spec(spec)
    sys.modules[name] = module
    spec.loader.exec_module(module)
    return module


@pytest.fixture(scope="module")
def kit():
    return _load("reportkit", KIT_PATH)


@pytest.fixture
def frame():
    return pd.DataFrame(
        {
            "team": ["Alpha", "Beta", "Gamma", "Delta"],
            "points": [120.4, 98.1, 77.7, 140.2],
            "wins": [9, 6, 4, 11],
        }
    )


def _exercise_pdf(kit, tmp_path, frame, theme):
    doc = kit.PdfReport("League Audit", "2026 season", theme=theme, accent="#c0392b")
    doc.title_page(classification="INTERNAL")
    doc.section("Executive Summary")
    doc.kpi_row([("Teams", 4), ("Avg points", 109.1), ("Leader", "Delta", "+30.1 avg")])
    doc.prose("Alpha is a **point-storage facility**.\n\nSecond paragraph with *italics*. <xml>")
    doc.bullets(["First finding", "Second finding"])
    doc.section("Standings")
    doc.table(frame, title="Full standings")
    doc.table([{"a": 1, "b": "x"}, {"a": 2, "b": "y"}])
    doc.table([[1, 2], [3, 4]], columns=["n1", "n2"])
    doc.callout("Hall of Fame", "Nobody. Everyone scored.")
    doc.chart(frame["team"], frame["points"], kind="bar", title="Points by team", ylabel="pts")
    doc.chart(frame["points"], kind="hist", title="Distribution")
    doc.chart([1, 2, 3], [4, 5, 3], kind="line", title="Trend")
    doc.chart([1, 2, 3], [4, 5, 3], kind="scatter")
    fig, ax = plt.subplots(figsize=(6, 3))
    ax.plot([1, 2, 3], [4, 3, 5])
    doc.figure(fig, caption="Custom matplotlib figure")
    plt.close(fig)
    doc.page_break()
    doc.section("Appendix")
    doc.prose("Methodology paragraph.")
    doc.table(frame, max_rows=2)
    out = doc.build(tmp_path / f"report-{theme}.pdf")
    return Path(out)


def test_pdf_report_builds_all_blocks_light_theme(kit, tmp_path, frame):
    out = _exercise_pdf(kit, tmp_path, frame, "light")
    data = out.read_bytes()
    assert data.startswith(b"%PDF-")
    assert len(data) > 20_000


def test_pdf_report_builds_dark_theme(kit, tmp_path, frame):
    out = _exercise_pdf(kit, tmp_path, frame, "dark")
    assert out.read_bytes().startswith(b"%PDF-")


def test_pdf_title_page_escapes_inputs_but_keeps_separator(kit, tmp_path):
    doc = kit.PdfReport("T & C")
    doc.title_page(classification="INTERNAL <draft>")
    joined = "\n".join(getattr(flowable, "text", "") for flowable in doc.story)
    assert "T &amp; C" in joined
    assert "INTERNAL &lt;draft&gt;" in joined
    assert "&middot;" in joined  # the separator entity must survive escaping
    assert "&amp;middot;" not in joined


def test_pdf_unknown_chart_kind_raises(kit, tmp_path):
    doc = kit.PdfReport("t")
    with pytest.raises(ValueError, match="unknown chart kind"):
        doc.chart([1, 2, 3], kind="donut")


def test_html_page_contains_all_blocks(kit, tmp_path, frame):
    page = kit.HtmlPage("League Dashboard", subtitle="live", theme="dark", accent="#c0392b")
    page.hero(kpis=[("Teams", 4), ("Avg points", 109.1)])
    page.section("Standings")
    page.table(frame)
    page.prose("Some **analysis** with <html> escaped.")
    page.bullets(["one", "two"])
    page.callout("Note", "A highlighted callout.")
    page.chart(frame["team"], frame["points"], kind="bar", title="Points")
    fig, ax = plt.subplots()
    ax.plot([1, 2], [3, 4])
    page.figure(fig, caption="cap")
    plt.close(fig)
    page.raw('<div class="custom">raw block</div>')
    out = Path(page.save(tmp_path / "page.html"))

    text = out.read_text(encoding="utf-8")
    assert text.startswith("<!DOCTYPE html>")
    for needle in [
        "League Dashboard",
        "Teams",
        "Alpha",
        "<b>analysis</b>",
        "&lt;html&gt;",
        "data:image/png;base64",
        "raw block",
        "cap",
    ]:
        assert needle in text, needle


def test_html_table_truncation_note(kit, tmp_path, frame):
    page = kit.HtmlPage("t")
    page.table(frame, max_rows=2)
    out = Path(page.save(tmp_path / "t.html"))
    text = out.read_text(encoding="utf-8")
    assert "2 more rows not shown" in text
    assert "Gamma" not in text


def test_seed_reportkit_copies_kit_into_workspace(tmp_path, monkeypatch):
    monkeypatch.setenv("SANDBOX_WORKSPACE", str(tmp_path))
    server = _load("sandboxd.server", SERVER_PATH)

    target = server.seed_reportkit(tmp_path / "ws")

    assert target == tmp_path / "ws" / "reportkit.py"
    assert target.read_bytes() == KIT_PATH.read_bytes()


def _slide_texts(slides):
    return [
        run.text
        for slide in slides
        for shape in slide.shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
    ]


def _exercise_deck(kit, tmp_path, frame, theme):
    deck = kit.Deck("League Audit", "2026 season", theme=theme, accent="#c0392b")
    deck.title_slide(classification="INTERNAL")
    deck.section("Findings")
    deck.slide("Overview")
    deck.kpi_row([("Teams", 4), ("Avg points", 109.1), ("Leader", "Delta", "+30.1 avg")])
    deck.prose("Alpha is a **point-storage facility**.\n\nSecond paragraph with *italics*.")
    deck.bullets(["First finding", "Second finding"])
    deck.table(frame, title="Full standings", max_rows=3)
    deck.callout("Hall of Fame", "Nobody. Everyone scored.")
    deck.chart(frame["team"], frame["points"], kind="bar", title="Points by team", ylabel="pts")
    fig, ax = plt.subplots(figsize=(6, 3))
    ax.plot([1, 2, 3], [4, 3, 5])
    deck.figure(fig, caption="Custom matplotlib figure")
    plt.close(fig)
    return Path(deck.build(tmp_path / f"deck-{theme}.pptx"))


def test_deck_builds_valid_pptx(kit, tmp_path, frame):
    from pptx import Presentation

    out = _exercise_deck(kit, tmp_path, frame, "light")
    assert out.read_bytes().startswith(b"PK")
    prs = Presentation(str(out))
    slides = list(prs.slides)
    assert len(slides) == 3
    assert prs.slide_width > prs.slide_height  # 16:9 landscape
    texts = _slide_texts(slides)
    assert any("League Audit" in text for text in texts)
    assert any("First finding" in text for text in texts)
    # Markdown-lite emphasis is stripped on slides, not rendered literally.
    assert any(text == "Alpha is a point-storage facility." for text in texts)


def test_deck_builds_dark_theme(kit, tmp_path, frame):
    out = _exercise_deck(kit, tmp_path, frame, "dark")
    assert out.read_bytes().startswith(b"PK")


def test_deck_content_requires_current_slide(kit):
    deck = kit.Deck("t")
    with pytest.raises(ValueError, match="slide"):
        deck.bullets(["no slide yet"])


def test_deck_table_truncates_and_figure_embeds_picture(kit, tmp_path, frame):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    deck = kit.Deck("t")
    deck.slide("Data")
    deck.table(frame, max_rows=2)
    fig, ax = plt.subplots()
    ax.bar(frame["team"], frame["points"])
    deck.figure(fig)
    plt.close(fig)
    out = deck.build(tmp_path / "d.pptx")

    slide = list(Presentation(str(out)).slides)[0]
    assert any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes)
    table = next(shape.table for shape in slide.shapes if shape.has_table)
    assert len(table.rows) == 3  # header plus the two kept rows
    texts = _slide_texts([slide])
    assert any("2 more rows not shown" in text for text in texts)
