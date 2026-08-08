# ruff: noqa: I001, E402
"""reportkit - compact report builders for Data Copilot sandbox sessions.

The agent decides everything about the document: sections, order, length,
topics, copy, charts, tables, and theming. This module absorbs the rendering
boilerplate (reportlab platypus plumbing, matplotlib chart styling, HTML/CSS
chrome) so authoring a report takes a few expressive lines instead of a wall
of low-level layout code. Raw reportlab / matplotlib / hand-written HTML
remain available for anything the builders cannot express.

Typical PDF use::

    import reportkit

    doc = reportkit.PdfReport("League Audit", "2026 season", accent="#c0392b")
    doc.title_page(classification="INTERNAL")
    doc.kpi_row([("Teams", 12), ("Avg points", "118.4")])
    doc.section("Standings")
    doc.table(standings_df)
    doc.section("Roasts")
    doc.prose("First paragraph.\\n\\nSecond paragraph with **bold** text.")
    doc.chart(labels, values, kind="bar", title="Points by team")
    doc.figure(my_matplotlib_fig)  # escape hatch: any custom chart
    doc.page_break()
    doc.section("Appendix")
    doc.prose(methodology_text)
    path = doc.build(Path(WORKSPACE) / "report.pdf")

HTML pages mirror the same API via ``HtmlPage(title, ...).save(path)``.
An appendix is just more sections; the agent controls document length.

PowerPoint decks follow the same philosophy via ``Deck(title, ...).build(path)``:
``title_slide()`` and ``section()`` are self-contained, ``slide(title)`` starts a
content slide, and the content methods (kpi_row, prose, bullets, table, callout,
chart, figure) append to it top to bottom.
"""

from __future__ import annotations

import base64
import html as _html_module
import io
import re
from collections.abc import Sequence
from datetime import UTC, datetime
from pathlib import Path
from typing import Any

import matplotlib

matplotlib.use("Agg")  # sandbox has no display; render charts headlessly

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import Image, PageBreak, Paragraph, Spacer
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle

_THEMES: dict[str, dict[str, str]] = {
    "light": {
        "bg": "#ffffff",
        "ink": "#1c2430",
        "muted": "#5b6672",
        "rule": "#d9dee6",
        "panel": "#f2f5fa",
    },
    "dark": {
        "bg": "#141a24",
        "ink": "#e8edf5",
        "muted": "#9aa7b8",
        "rule": "#2c3646",
        "panel": "#1d2735",
    },
}

_BOLD_RE = re.compile(r"\*\*(.+?)\*\*")
_ITALIC_RE = re.compile(r"\*(.+?)\*")


def _palette(theme: str, accent: str) -> dict[str, str]:
    pal = dict(_THEMES.get(theme, _THEMES["light"]))
    pal["accent"] = accent
    return pal


def _md_lite(text: str) -> str:
    """Escape XML, then honor a minimal markdown subset (**bold**, *italic*)."""
    escaped = _html_module.escape(text)
    escaped = _BOLD_RE.sub(r"<b>\1</b>", escaped)
    return _ITALIC_RE.sub(r"<i>\1</i>", escaped)


def _plain(text: str) -> str:
    """Strip markdown-lite emphasis for plain PowerPoint text runs."""
    text = _BOLD_RE.sub(r"\1", str(text))
    return _ITALIC_RE.sub(r"\1", text)


def _flow_height(text: str, *, size: float, width: float) -> float:
    """Rough rendered height (inches) for wrapped text in a slide textbox."""
    chars = max(20, int(width * 10.5 * 14 / size))
    lines = max(1, -(-len(text) // chars))
    return lines * size * 1.5 / 72


def _hex(value: str) -> colors.Color:
    return colors.HexColor(value)


def _tint(hex_color: str, factor: float = 0.82) -> colors.Color:
    """Blend a hex color toward white; used for soft header/background fills."""
    base = _hex(hex_color)
    white = colors.Color(1, 1, 1)
    return colors.Color(
        base.red + (white.red - base.red) * factor,
        base.green + (white.green - base.green) * factor,
        base.blue + (white.blue - base.blue) * factor,
    )


def _tint_css(hex_color: str, factor: float = 0.82) -> str:
    tint = _tint(hex_color, factor)
    return f"#{int(tint.red * 255):02x}{int(tint.green * 255):02x}{int(tint.blue * 255):02x}"


def _as_list(values: Any) -> list[Any]:
    """Accept lists, tuples, numpy arrays, and pandas Series uniformly."""
    if hasattr(values, "tolist"):
        return list(values.tolist())
    return list(values)


def _rows_from(
    data: Any, columns: Sequence[str] | None, max_rows: int
) -> tuple[list[str], list[list[str]], int]:
    """Normalize a DataFrame / list-of-dicts / list-of-lists into string rows.

    Returns (header, rows, hidden_row_count).
    """
    if hasattr(data, "to_dict") and hasattr(data, "columns"):  # pandas DataFrame
        header = [str(c) for c in data.columns] if columns is None else [str(c) for c in columns]
        frame = data if columns is None else data[list(columns)]
        records = frame.to_dict("records")
        rows = [[_cell(record.get(col)) for col in header] for record in records]
    elif data and isinstance(data[0], dict):
        header = [str(c) for c in (columns or list(data[0].keys()))]
        rows = [[_cell(record.get(col)) for col in header] for record in data]
    else:
        matrix = [list(row) for row in data]
        width = max((len(row) for row in matrix), default=0)
        header = [str(c) for c in (columns or [f"col {i + 1}" for i in range(width)])]
        rows = [[_cell(value) for value in row] for row in matrix]
    hidden = max(0, len(rows) - max_rows)
    return header, rows[:max_rows], hidden


def _cell(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, float):
        return f"{value:,.3g}" if abs(value) < 1e6 else f"{value:,.0f}"
    if isinstance(value, int):
        return f"{value:,}"
    return str(value)


_CHART_FIGSIZE = (7.0, 3.4)


def _chart_png(
    x: Any,
    y: Any,
    *,
    kind: str,
    title: str,
    xlabel: str,
    ylabel: str,
    pal: dict[str, str],
) -> bytes:
    """Render a quick themed chart to PNG bytes.

    ``x``/``y`` take lists, numpy arrays, or pandas Series. ``kind`` is one of
    bar, line, scatter, or hist (hist only needs ``x``).
    """
    fig, ax = plt.subplots(figsize=_CHART_FIGSIZE, dpi=160)
    _style_axes(fig, ax, pal)
    xs = _as_list(x)
    ys = None if y is None else _as_list(y)
    if kind == "hist":
        ax.hist(xs, bins=min(30, max(5, len(xs) // 8)), color=pal["accent"], edgecolor=pal["bg"])
    elif kind == "scatter":
        if ys is None:
            raise ValueError("scatter charts need both x and y")
        ax.scatter(xs, ys, color=pal["accent"], alpha=0.75, s=28, edgecolors="none")
    elif kind == "line":
        if ys is None:
            ys, xs = xs, list(range(len(xs)))
        ax.plot(xs, ys, color=pal["accent"], linewidth=2, marker="o", markersize=4)
    elif kind == "bar":
        if ys is None:
            ys, xs = xs, [str(i + 1) for i in range(len(xs))]
        ax.bar([str(v) for v in xs], ys, color=pal["accent"])
        if len(xs) > 8:
            plt.setp(ax.get_xticklabels(), rotation=45, ha="right")
    else:
        plt.close(fig)
        raise ValueError(f"unknown chart kind: {kind!r} (expected bar, line, scatter, or hist)")
    if title:
        ax.set_title(title, color=pal["ink"], fontsize=12, fontweight="bold", pad=10)
    if xlabel:
        ax.set_xlabel(xlabel, color=pal["muted"], fontsize=9)
    if ylabel:
        ax.set_ylabel(ylabel, color=pal["muted"], fontsize=9)
    fig.tight_layout()
    png = _fig_png(fig)
    plt.close(fig)
    return png


def _style_axes(fig: Any, ax: Any, pal: dict[str, str]) -> None:
    fig.patch.set_facecolor(pal["bg"])
    ax.set_facecolor(pal["bg"])
    for spine in ("top", "right"):
        ax.spines[spine].set_visible(False)
    for spine in ("left", "bottom"):
        ax.spines[spine].set_color(pal["rule"])
    ax.tick_params(colors=pal["muted"], labelsize=8)
    ax.grid(axis="y", color=pal["rule"], linewidth=0.6, alpha=0.7)
    ax.set_axisbelow(True)


def _fig_png(fig: Any) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=160, bbox_inches="tight", facecolor=fig.get_facecolor())
    return buf.getvalue()


def _flowable_image(png: bytes, width: float) -> Image:
    """Size an Image flowable from its PNG bytes, preserving aspect ratio."""
    # PNG header carries width/height big-endian at bytes 16..24.
    px_w = int.from_bytes(png[16:20], "big")
    px_h = int.from_bytes(png[20:24], "big")
    height = width * (px_h / px_w) if px_w else width * 0.5
    return Image(io.BytesIO(png), width=width, height=height)


def _img_tag(png: bytes, alt: str) -> str:
    encoded = base64.b64encode(png).decode("ascii")
    return f'<img src="data:image/png;base64,{encoded}" alt="{_html_module.escape(alt)}"/>'


def _today() -> str:
    return datetime.now(UTC).strftime("%B %d, %Y")


class PdfReport:
    """Compose a styled multi-page PDF from sections, tables, and charts.

    The caller controls structure, length, topics, copy, and theming; this
    class only handles rendering mechanics.
    """

    def __init__(
        self,
        title: str,
        subtitle: str = "",
        *,
        theme: str = "light",
        accent: str = "#2f6fed",
        author: str = "Data Copilot",
    ) -> None:
        self.title = title
        self.subtitle = subtitle
        self.author = author
        self.pal = _palette(theme, accent)
        self.story: list[Any] = []
        ink, muted = _hex(self.pal["ink"]), _hex(self.pal["muted"])
        self.styles = {
            "title": ParagraphStyle(
                "rk-title", fontName="Helvetica-Bold", fontSize=30, leading=36,
                textColor=ink, alignment=1, spaceAfter=10,
            ),
            "subtitle": ParagraphStyle(
                "rk-subtitle", fontName="Helvetica", fontSize=14, leading=19,
                textColor=muted, alignment=1, spaceAfter=6,
            ),
            "meta": ParagraphStyle(
                "rk-meta", fontName="Helvetica", fontSize=10, leading=14,
                textColor=muted, alignment=1,
            ),
            "h1": ParagraphStyle(
                "rk-h1", fontName="Helvetica-Bold", fontSize=17, leading=21,
                textColor=_hex(self.pal["accent"]), spaceBefore=16, spaceAfter=7,
            ),
            "h2": ParagraphStyle(
                "rk-h2", fontName="Helvetica-Bold", fontSize=13, leading=17,
                textColor=ink, spaceBefore=11, spaceAfter=5,
            ),
            "body": ParagraphStyle(
                "rk-body", fontName="Helvetica", fontSize=10.5, leading=15,
                textColor=ink, spaceAfter=7,
            ),
            "bullet": ParagraphStyle(
                "rk-bullet", fontName="Helvetica", fontSize=10.5, leading=15,
                textColor=ink, leftIndent=16, bulletIndent=6, spaceAfter=3,
            ),
            "caption": ParagraphStyle(
                "rk-caption", fontName="Helvetica-Oblique", fontSize=9, leading=12,
                textColor=muted, spaceBefore=3, spaceAfter=8,
            ),
        }

    def title_page(self, *, classification: str | None = None, date: str | None = None) -> None:
        """Add a cover page (followed by a page break) with title block."""
        self.story.append(Spacer(1, 2.1 * inch))
        rule = Table([[""]], colWidths=[2.4 * inch], rowHeights=[3])
        rule.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, -1), _hex(self.pal["accent"]))]))
        rule.hAlign = "CENTER"
        self.story.append(rule)
        self.story.append(Spacer(1, 0.35 * inch))
        self.story.append(Paragraph(_md_lite(self.title), self.styles["title"]))
        if self.subtitle:
            self.story.append(Paragraph(_md_lite(self.subtitle), self.styles["subtitle"]))
        self.story.append(Spacer(1, 0.5 * inch))
        meta = _md_lite(date or _today())
        if classification:
            meta = f"{_md_lite(classification)} &middot; {meta}"
        self.story.append(Paragraph(meta, self.styles["meta"]))
        self.story.append(Paragraph(_md_lite(f"Prepared by {self.author}"), self.styles["meta"]))
        self.page_break()

    def section(self, title: str, *, level: int = 1) -> None:
        self.story.append(Paragraph(_md_lite(title), self.styles["h1" if level == 1 else "h2"]))

    def prose(self, text: str) -> None:
        """Add paragraphs; blank lines separate them, **bold**/*italic* honored."""
        for paragraph in text.strip().split("\n\n"):
            if paragraph.strip():
                self.story.append(Paragraph(_md_lite(paragraph.strip()), self.styles["body"]))

    def bullets(self, items: Sequence[Any]) -> None:
        for item in items:
            self.story.append(
                Paragraph(_md_lite(str(item)), self.styles["bullet"], bulletText="\u2022")
            )

    def kpi_row(self, kpis: Sequence[Sequence[Any]]) -> None:
        """Add a row of KPI tiles; each entry is (label, value[, note])."""
        cells = []
        for entry in kpis:
            label, value = str(entry[0]), _cell(entry[1])
            note = str(entry[2]) if len(entry) > 2 else ""
            tile_style = ParagraphStyle(
                "rk-kpi", parent=self.styles["body"], alignment=1, spaceAfter=0
            )
            parts = [
                Paragraph(f'<font size="16"><b>{_md_lite(value)}</b></font>', tile_style),
                Paragraph(
                    f'<font size="8" color="{self.pal["muted"]}">{_md_lite(label.upper())}</font>',
                    tile_style,
                ),
            ]
            if note:
                parts.append(
                    Paragraph(
                        f'<font size="7" color="{self.pal["muted"]}">{_md_lite(note)}</font>',
                        tile_style,
                    )
                )
            cells.append(parts)
        if not cells:
            return
        width = 6.9 * inch / len(cells)
        table = Table([cells], colWidths=[width] * len(cells))
        table.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, -1), _hex(self.pal["panel"])),
                    ("BOX", (0, 0), (-1, -1), 0.75, _hex(self.pal["rule"])),
                    ("INNERGRID", (0, 0), (-1, -1), 0.75, _hex(self.pal["rule"])),
                    ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                    ("TOPPADDING", (0, 0), (-1, -1), 10),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 10),
                ]
            )
        )
        self.story.append(Spacer(1, 0.06 * inch))
        self.story.append(table)
        self.story.append(Spacer(1, 0.12 * inch))

    def table(
        self,
        data: Any,
        *,
        columns: Sequence[str] | None = None,
        title: str | None = None,
        max_rows: int = 40,
    ) -> None:
        """Add a styled table from a DataFrame, list of dicts, or list of rows."""
        header, rows, hidden = _rows_from(data, columns, max_rows)
        if not header:
            return
        if title:
            self.story.append(Paragraph(_md_lite(title), self.styles["h2"]))
        body = [
            [Paragraph(_md_lite(cell), self.styles["body"]) for cell in row]
            for row in [header, *rows]
        ]
        count = len(header)
        col_width = 6.9 * inch / max(count, 1)
        table = Table(body, colWidths=[col_width] * count, repeatRows=1)
        table.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), _tint(self.pal["accent"])),
                    ("ROWBACKGROUNDS", (0, 1), (-1, -1), [_hex(self.pal["bg"]), _hex(self.pal["panel"])]),
                    ("LINEBELOW", (0, 0), (-1, 0), 1, _hex(self.pal["accent"])),
                    ("GRID", (0, 0), (-1, -1), 0.5, _hex(self.pal["rule"])),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("TOPPADDING", (0, 0), (-1, -1), 4),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                    ("LEFTPADDING", (0, 0), (-1, -1), 6),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                ]
            )
        )
        self.story.append(table)
        if hidden:
            self.story.append(
                Paragraph(f"... {hidden} more rows not shown.", self.styles["caption"])
            )
        self.story.append(Spacer(1, 0.1 * inch))

    def callout(self, title: str, text: str) -> None:
        """Add an accent-bordered highlight box."""
        inner = [
            Paragraph(
                f'<font color="{self.pal["accent"]}"><b>{_md_lite(title)}</b></font>',
                self.styles["body"],
            )
        ]
        for paragraph in text.strip().split("\n\n"):
            if paragraph.strip():
                inner.append(Paragraph(_md_lite(paragraph.strip()), self.styles["body"]))
        box = Table([[inner]], colWidths=[6.9 * inch])
        box.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, -1), _hex(self.pal["panel"])),
                    ("LINEBEFORE", (0, 0), (0, -1), 3, _hex(self.pal["accent"])),
                    ("TOPPADDING", (0, 0), (-1, -1), 8),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
                    ("LEFTPADDING", (0, 0), (-1, -1), 10),
                ]
            )
        )
        self.story.append(box)
        self.story.append(Spacer(1, 0.1 * inch))

    def chart(
        self,
        x: Any,
        y: Any = None,
        *,
        kind: str = "bar",
        title: str = "",
        xlabel: str = "",
        ylabel: str = "",
        width: float = 6.3 * inch,
    ) -> None:
        """Add a quick themed chart (bar/line/scatter/hist) straight from data."""
        png = _chart_png(
            x, y, kind=kind, title=title, xlabel=xlabel, ylabel=ylabel, pal=self.pal
        )
        image = _flowable_image(png, width)
        image.hAlign = "CENTER"
        self.story.append(image)
        self.story.append(Spacer(1, 0.08 * inch))

    def figure(self, fig: Any, *, caption: str | None = None, width: float = 6.3 * inch) -> None:
        """Embed any matplotlib figure for full chart control."""
        image = _flowable_image(_fig_png(fig), width)
        image.hAlign = "CENTER"
        self.story.append(image)
        if caption:
            self.story.append(Paragraph(_md_lite(caption), self.styles["caption"]))
        else:
            self.story.append(Spacer(1, 0.08 * inch))

    def page_break(self) -> None:
        self.story.append(PageBreak())

    def spacer(self, height: float = 0.15 * inch) -> None:
        self.story.append(Spacer(1, height))

    def build(self, path: Any) -> str:
        """Render the PDF to ``path`` and return it as a string."""
        target = Path(path)
        target.parent.mkdir(parents=True, exist_ok=True)
        doc = SimpleDocTemplate(
            str(target),
            pagesize=letter,
            leftMargin=0.8 * inch,
            rightMargin=0.8 * inch,
            topMargin=0.8 * inch,
            bottomMargin=0.8 * inch,
            title=self.title,
            author=self.author,
        )
        pal = self.pal
        title = self.title

        def decorate(canvas: Any, _: Any) -> None:
            canvas.saveState()
            if pal["bg"] != "#ffffff":
                canvas.setFillColor(_hex(pal["bg"]))
                canvas.rect(0, 0, letter[0], letter[1], stroke=0, fill=1)
            canvas.setStrokeColor(_hex(pal["rule"]))
            canvas.setLineWidth(0.5)
            canvas.line(0.8 * inch, 0.62 * inch, letter[0] - 0.8 * inch, 0.62 * inch)
            canvas.setFillColor(_hex(pal["muted"]))
            canvas.setFont("Helvetica", 8)
            canvas.drawString(0.8 * inch, 0.45 * inch, title[:80])
            canvas.drawRightString(letter[0] - 0.8 * inch, 0.45 * inch, f"Page {canvas.getPageNumber()}")
            canvas.restoreState()

        doc.build(self.story, onFirstPage=decorate, onLaterPages=decorate)
        return str(target)


class HtmlPage:
    """Compose a styled single-file HTML page (charts embedded as base64)."""

    def __init__(
        self,
        title: str,
        *,
        subtitle: str = "",
        theme: str = "light",
        accent: str = "#2f6fed",
    ) -> None:
        self.title = title
        self.subtitle = subtitle
        self.pal = _palette(theme, accent)
        self.parts: list[str] = []

    def hero(self, subtitle: str = "", kpis: Sequence[Sequence[Any]] | None = None) -> None:
        """Add a header banner with title, subtitle, and optional KPI tiles."""
        sub = subtitle or self.subtitle
        self.parts.append('<header class="hero">')
        self.parts.append(f"<h1>{_md_lite(self.title)}</h1>")
        if sub:
            self.parts.append(f'<p class="subtitle">{_md_lite(sub)}</p>')
        self.parts.append(f'<p class="meta">{_md_lite(_today())}</p>')
        self.parts.append("</header>")
        if kpis:
            self.kpi_row(kpis)

    def section(self, title: str) -> None:
        self.parts.append(f"<h2>{_md_lite(title)}</h2>")

    def prose(self, text: str) -> None:
        for paragraph in text.strip().split("\n\n"):
            if paragraph.strip():
                self.parts.append(f"<p>{_md_lite(paragraph.strip())}</p>")

    def bullets(self, items: Sequence[Any]) -> None:
        self.parts.append("<ul>")
        for item in items:
            self.parts.append(f"<li>{_md_lite(str(item))}</li>")
        self.parts.append("</ul>")

    def kpi_row(self, kpis: Sequence[Sequence[Any]]) -> None:
        self.parts.append('<div class="kpis">')
        for entry in kpis:
            label, value = str(entry[0]), _cell(entry[1])
            note = f'<span class="note">{_md_lite(str(entry[2]))}</span>' if len(entry) > 2 else ""
            self.parts.append(
                f'<div class="kpi"><span class="value">{_md_lite(value)}</span>'
                f'<span class="label">{_md_lite(label)}</span>{note}</div>'
            )
        self.parts.append("</div>")

    def table(
        self, data: Any, *, columns: Sequence[str] | None = None, max_rows: int = 100
    ) -> None:
        header, rows, hidden = _rows_from(data, columns, max_rows)
        if not header:
            return
        self.parts.append("<table><thead><tr>")
        self.parts.extend(f"<th>{_md_lite(col)}</th>" for col in header)
        self.parts.append("</tr></thead><tbody>")
        for row in rows:
            self.parts.append("<tr>" + "".join(f"<td>{_md_lite(cell)}</td>" for cell in row) + "</tr>")
        self.parts.append("</tbody></table>")
        if hidden:
            self.parts.append(f'<p class="caption">... {hidden} more rows not shown.</p>')

    def callout(self, title: str, text: str) -> None:
        body = "".join(
            f"<p>{_md_lite(p.strip())}</p>" for p in text.strip().split("\n\n") if p.strip()
        )
        self.parts.append(
            f'<div class="callout"><strong>{_md_lite(title)}</strong>{body}</div>'
        )

    def chart(
        self,
        x: Any,
        y: Any = None,
        *,
        kind: str = "bar",
        title: str = "",
        xlabel: str = "",
        ylabel: str = "",
    ) -> None:
        png = _chart_png(
            x, y, kind=kind, title=title, xlabel=xlabel, ylabel=ylabel, pal=self.pal
        )
        self.parts.append(f'<figure>{_img_tag(png, title or kind)}</figure>')

    def figure(self, fig: Any, *, caption: str | None = None) -> None:
        tag = _img_tag(_fig_png(fig), caption or "chart")
        cap = f"<figcaption>{_md_lite(caption)}</figcaption>" if caption else ""
        self.parts.append(f"<figure>{tag}{cap}</figure>")

    def raw(self, html: str) -> None:
        """Insert raw markup for anything the helpers cannot express."""
        self.parts.append(html)

    def save(self, path: Any) -> str:
        target = Path(path)
        target.parent.mkdir(parents=True, exist_ok=True)
        pal = self.pal
        css = f"""
:root {{ --accent: {pal['accent']}; }}
* {{ box-sizing: border-box; }}
body {{ margin: 0; background: {pal['bg']}; color: {pal['ink']};
  font-family: -apple-system, 'Segoe UI', Roboto, Helvetica, Arial, sans-serif;
  line-height: 1.55; }}
main {{ max-width: 860px; margin: 0 auto; padding: 28px 22px 64px; }}
.hero {{ border-bottom: 3px solid var(--accent); padding: 18px 0 14px; margin-bottom: 18px; }}
.hero h1 {{ margin: 0 0 4px; font-size: 30px; }}
.subtitle {{ margin: 0; color: {pal['muted']}; font-size: 16px; }}
.meta {{ margin: 6px 0 0; color: {pal['muted']}; font-size: 12px; }}
h2 {{ color: var(--accent); margin: 28px 0 10px; font-size: 21px; }}
p {{ margin: 0 0 10px; }}
ul {{ margin: 0 0 12px; padding-left: 22px; }}
.kpis {{ display: flex; flex-wrap: wrap; gap: 10px; margin: 14px 0 6px; }}
.kpi {{ flex: 1 1 120px; background: {pal['panel']}; border: 1px solid {pal['rule']};
  border-radius: 8px; padding: 12px 14px; text-align: center; }}
.kpi .value {{ display: block; font-size: 22px; font-weight: 700; }}
.kpi .label {{ display: block; font-size: 11px; text-transform: uppercase;
  letter-spacing: 0.04em; color: {pal['muted']}; margin-top: 2px; }}
.kpi .note {{ display: block; font-size: 11px; color: {pal['muted']}; margin-top: 2px; }}
table {{ border-collapse: collapse; width: 100%; margin: 10px 0 16px; font-size: 14px; }}
th {{ background: {_tint_css(pal['accent'])}; text-align: left; }}
th, td {{ border: 1px solid {pal['rule']}; padding: 6px 9px; }}
tbody tr:nth-child(even) {{ background: {pal['panel']}; }}
.callout {{ background: {pal['panel']}; border-left: 4px solid var(--accent);
  border-radius: 0 8px 8px 0; padding: 10px 14px; margin: 12px 0; }}
figure {{ margin: 14px 0; text-align: center; }}
figure img {{ max-width: 100%; border: 1px solid {pal['rule']}; border-radius: 8px; }}
figcaption, .caption {{ color: {pal['muted']}; font-size: 12px; margin-top: 4px; }}
"""
        document = (
            "<!DOCTYPE html>\n<html lang=\"en\"><head><meta charset=\"utf-8\"/>"
            '<meta name="viewport" content="width=device-width, initial-scale=1"/>'
            f"<title>{_html_module.escape(self.title)}</title><style>{css}</style></head>"
            f"<body><main>{''.join(self.parts)}</main></body></html>"
        )
        target.write_text(document, encoding="utf-8")
        return str(target)


class Deck:
    """Compose a styled 16:9 PowerPoint deck (.pptx) from slides and charts.

    Same philosophy as PdfReport: the caller decides structure, copy, and
    theming; this class absorbs the python-pptx mechanics. ``title_slide()``
    and ``section()`` are self-contained slides; ``slide(title)`` starts a
    content slide that the content methods (prose, bullets, kpi_row, table,
    callout, chart, figure) append to, flowing top to bottom.
    """

    _SLIDE_W = 13.333
    _SLIDE_H = 7.5
    _MARGIN = 0.6
    _CONTENT_W = _SLIDE_W - 2 * _MARGIN

    def __init__(
        self,
        title: str,
        subtitle: str = "",
        *,
        theme: str = "light",
        accent: str = "#2f6fed",
        author: str = "Data Copilot",
    ) -> None:
        self.title = title
        self.subtitle = subtitle
        self.author = author
        self.pal = _palette(theme, accent)
        self._prs = Presentation()
        self._prs.slide_width = Inches(self._SLIDE_W)
        self._prs.slide_height = Inches(self._SLIDE_H)
        self._slide: Any = None
        self._cursor = 0.0

    @staticmethod
    def _rgb(hex_color: str) -> RGBColor:
        return RGBColor.from_string(hex_color.lstrip("#").upper())

    def _blank_slide(self) -> Any:
        slide = self._prs.slides.add_slide(self._prs.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self._rgb(self.pal["bg"])
        return slide

    @staticmethod
    def _textbox(slide: Any, left: float, top: float, width: float, height: float) -> Any:
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        box.text_frame.word_wrap = True
        return box.text_frame

    def _para(
        self,
        tf: Any,
        text: str,
        *,
        size: float,
        color: str,
        bold: bool = False,
        align: Any = PP_ALIGN.LEFT,
        space_after: float = 4,
        first: bool = False,
    ) -> None:
        paragraph = tf.paragraphs[0] if first else tf.add_paragraph()
        paragraph.alignment = align
        paragraph.space_after = Pt(space_after)
        run = paragraph.add_run()
        run.text = _plain(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = self._rgb(color)

    def _bar(self, slide: Any, left: float, top: float, width: float, height: float, color: str) -> None:
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._rgb(color)
        bar.line.fill.background()
        bar.shadow.inherit = False

    def _current(self) -> Any:
        if self._slide is None:
            raise ValueError("no current slide: call .slide(title) before adding content")
        return self._slide

    def title_slide(self, *, classification: str | None = None, date: str | None = None) -> None:
        """Add a cover slide with the deck title block."""
        slide = self._blank_slide()
        self._bar(slide, (self._SLIDE_W - 2.4) / 2, 2.02, 2.4, 0.045, self.pal["accent"])
        tf = self._textbox(slide, 1.0, 2.3, self._SLIDE_W - 2.0, 1.6)
        self._para(
            tf, self.title, size=40, color=self.pal["ink"], bold=True,
            align=PP_ALIGN.CENTER, first=True,
        )
        if self.subtitle:
            self._para(tf, self.subtitle, size=18, color=self.pal["muted"], align=PP_ALIGN.CENTER)
        meta = date or _today()
        if classification:
            meta = f"{classification} · {meta}"
        tf = self._textbox(slide, 1.0, 5.9, self._SLIDE_W - 2.0, 0.9)
        self._para(tf, meta, size=11, color=self.pal["muted"], align=PP_ALIGN.CENTER, first=True)
        self._para(tf, f"Prepared by {self.author}", size=11, color=self.pal["muted"], align=PP_ALIGN.CENTER)
        self._slide = None

    def section(self, title: str) -> None:
        """Add a section-divider slide with an accent side bar."""
        slide = self._blank_slide()
        self._bar(slide, 0.92, 3.02, 0.09, 1.2, self.pal["accent"])
        tf = self._textbox(slide, 1.25, 2.95, 10.8, 1.5)
        self._para(tf, title, size=32, color=self.pal["accent"], bold=True, first=True)
        self._slide = None

    def slide(self, title: str) -> None:
        """Start a content slide; later content methods append to it."""
        slide = self._blank_slide()
        tf = self._textbox(slide, self._MARGIN, 0.32, self._CONTENT_W, 0.75)
        self._para(tf, title, size=26, color=self.pal["ink"], bold=True, first=True)
        self._bar(slide, self._MARGIN + 0.02, 1.08, 1.6, 0.045, self.pal["accent"])
        self._slide = slide
        self._cursor = 1.32

    def prose(self, text: str) -> None:
        """Add paragraphs to the current slide (**bold**/*italic* are stripped)."""
        slide = self._current()
        paragraphs = [p.strip() for p in text.strip().split("\n\n") if p.strip()]
        height = sum(_flow_height(p, size=14, width=self._CONTENT_W) for p in paragraphs) + 0.2
        tf = self._textbox(slide, self._MARGIN, self._cursor, self._CONTENT_W, height)
        for index, paragraph in enumerate(paragraphs):
            self._para(tf, paragraph, size=14, color=self.pal["ink"], first=index == 0, space_after=6)
        self._cursor += height

    def bullets(self, items: Sequence[Any]) -> None:
        """Add bullet points to the current slide."""
        slide = self._current()
        entries = [str(item) for item in items]
        height = (
            sum(_flow_height(item, size=15, width=self._CONTENT_W - 0.3) for item in entries)
            + 0.08 * len(entries)
            + 0.15
        )
        tf = self._textbox(slide, self._MARGIN + 0.15, self._cursor, self._CONTENT_W - 0.15, height)
        for index, item in enumerate(entries):
            self._para(
                tf, f"• {item}", size=15, color=self.pal["ink"], first=index == 0, space_after=6
            )
        self._cursor += height

    def kpi_row(self, kpis: Sequence[Sequence[Any]]) -> None:
        """Add a row of KPI tiles to the current slide; each entry is (label, value[, note])."""
        slide = self._current()
        entries = list(kpis)
        if not entries:
            return
        gap = 0.2
        width = (self._CONTENT_W - gap * (len(entries) - 1)) / len(entries)
        for index, entry in enumerate(entries):
            left = self._MARGIN + index * (width + gap)
            tile = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(self._cursor), Inches(width), Inches(1.15),
            )
            tile.fill.solid()
            tile.fill.fore_color.rgb = self._rgb(self.pal["panel"])
            tile.line.color.rgb = self._rgb(self.pal["rule"])
            tile.line.width = Pt(0.75)
            tile.shadow.inherit = False
            tf = tile.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            self._para(
                tf, _cell(entry[1]), size=22, color=self.pal["accent"], bold=True,
                align=PP_ALIGN.CENTER, first=True, space_after=2,
            )
            self._para(
                tf, str(entry[0]).upper(), size=9, color=self.pal["muted"],
                align=PP_ALIGN.CENTER, space_after=0,
            )
            if len(entry) > 2:
                self._para(
                    tf, str(entry[2]), size=8, color=self.pal["muted"],
                    align=PP_ALIGN.CENTER, space_after=0,
                )
        self._cursor += 1.15 + 0.2

    def _style_cell(self, cell: Any, text: str, *, fill: str, size: float, bold: bool = False) -> None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = self._rgb(fill)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)
        cell.text_frame.word_wrap = True
        self._para(
            cell.text_frame, text, size=size, color=self.pal["ink"], bold=bold,
            first=True, space_after=0,
        )

    def table(
        self,
        data: Any,
        *,
        columns: Sequence[str] | None = None,
        title: str | None = None,
        max_rows: int = 10,
    ) -> None:
        """Add a styled table to the current slide (DataFrame / dicts / rows)."""
        slide = self._current()
        header, rows, hidden = _rows_from(data, columns, max_rows)
        if not header:
            return
        if title:
            tf = self._textbox(slide, self._MARGIN, self._cursor, self._CONTENT_W, 0.35)
            self._para(tf, title, size=14, color=self.pal["ink"], bold=True, first=True)
            self._cursor += 0.4
        count = len(header)
        row_height = 0.34
        height = row_height * (len(rows) + 1)
        frame = slide.shapes.add_table(
            len(rows) + 1, count,
            Inches(self._MARGIN), Inches(self._cursor), Inches(self._CONTENT_W), Inches(height),
        )
        table = frame.table
        table.first_row = False
        table.horz_banding = False
        for column in table.columns:
            column.width = Emu(int(Inches(self._CONTENT_W) / count))
        for row in table.rows:
            row.height = Emu(int(Inches(row_height)))
        tint = _tint_css(self.pal["accent"])
        for j, label in enumerate(header):
            self._style_cell(table.cell(0, j), label, fill=tint, size=11, bold=True)
        for i, row in enumerate(rows, 1):
            fill = self.pal["bg"] if i % 2 else self.pal["panel"]
            for j, value in enumerate(row):
                self._style_cell(table.cell(i, j), value, fill=fill, size=10.5)
        self._cursor += height + 0.15
        if hidden:
            tf = self._textbox(slide, self._MARGIN, self._cursor, self._CONTENT_W, 0.3)
            self._para(tf, f"... {hidden} more rows not shown.", size=10, color=self.pal["muted"], first=True)
            self._cursor += 0.32

    def callout(self, title: str, text: str) -> None:
        """Add an accent-bordered highlight box to the current slide."""
        slide = self._current()
        paragraphs = [p.strip() for p in text.strip().split("\n\n") if p.strip()]
        height = (
            0.45
            + sum(_flow_height(p, size=12.5, width=self._CONTENT_W - 0.4) for p in paragraphs)
            + 0.1 * len(paragraphs)
        )
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(self._MARGIN), Inches(self._cursor), Inches(self._CONTENT_W), Inches(height),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = self._rgb(self.pal["panel"])
        box.line.color.rgb = self._rgb(self.pal["accent"])
        box.line.width = Pt(1.5)
        box.shadow.inherit = False
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.1)
        self._para(tf, title, size=12.5, color=self.pal["accent"], bold=True, first=True, space_after=4)
        for paragraph in paragraphs:
            self._para(tf, paragraph, size=12.5, color=self.pal["ink"], space_after=4)
        self._cursor += height + 0.18

    def _place_image(self, slide: Any, png: bytes, width: float) -> None:
        px_w = int.from_bytes(png[16:20], "big")
        px_h = int.from_bytes(png[20:24], "big")
        height = width * (px_h / px_w) if px_w else width * 0.5
        left = (self._SLIDE_W - width) / 2
        slide.shapes.add_picture(
            io.BytesIO(png), Inches(left), Inches(self._cursor), width=Inches(width)
        )
        self._cursor += height + 0.18

    def chart(
        self,
        x: Any,
        y: Any = None,
        *,
        kind: str = "bar",
        title: str = "",
        xlabel: str = "",
        ylabel: str = "",
        width: float = 7.2,
    ) -> None:
        """Add a quick themed chart (bar/line/scatter/hist) to the current slide."""
        slide = self._current()
        png = _chart_png(x, y, kind=kind, title=title, xlabel=xlabel, ylabel=ylabel, pal=self.pal)
        self._place_image(slide, png, width)

    def figure(self, fig: Any, *, caption: str | None = None, width: float = 7.2) -> None:
        """Embed any matplotlib figure on the current slide."""
        slide = self._current()
        self._place_image(slide, _fig_png(fig), width)
        if caption:
            tf = self._textbox(slide, self._MARGIN, self._cursor, self._CONTENT_W, 0.3)
            self._para(
                tf, caption, size=10, color=self.pal["muted"], align=PP_ALIGN.CENTER, first=True
            )
            self._cursor += 0.3

    def build(self, path: Any) -> str:
        """Save the deck to ``path`` (.pptx) and return it as a string."""
        target = Path(path)
        target.parent.mkdir(parents=True, exist_ok=True)
        self._prs.save(str(target))
        return str(target)
